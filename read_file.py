#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
本地知识库 - 文件读取脚本 (重构版)
支持格式: docx, pdf, xlsx, pptx, txt, md
增强: PDF表格支持、语义切片
"""

import os
import re
import sys

def read_docx(file_path):
    """读取 Word 文档"""
    try:
        from docx import Document
        doc = Document(file_path)
        content = []
        for para in doc.paragraphs:
            if para.text.strip():
                content.append(para.text)
        for table in doc.tables:
            # 表格转为Markdown格式
            table_md = table_to_markdown(table)
            if table_md:
                content.append(table_md)
        return "\n".join(content)
    except Exception as e:
        return f"Error reading docx: {str(e)}"

def table_to_markdown(table):
    """将Word表格转为Markdown"""
    if not table.rows:
        return ""
    
    rows = []
    for i, row in enumerate(table.rows):
        cells = [cell.text.strip() for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
        if i == 0:
            rows.append("|" + "|".join(["---"] * len(cells)) + "|")
    
    return "\n".join(rows)

def read_pdf_with_tables(file_path):
    """使用pdfplumber读取PDF（含表格）"""
    try:
        import pdfplumber
        from pdfplumber.utils import extract_text, get_bbox_overlap, obj_to_bbox
        
        text_parts = []
        
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                # 提取表格
                tables = page.extract_tables()
                
                # 提取文本
                text = page.extract_text()
                
                if tables:
                    # 有表格时，优先保留表格结构
                    for table in tables:
                        if table:
                            table_md = markdown_table(table)
                            if table_md:
                                text_parts.append(f"\n[表格 {page_num}]\n")
                                text_parts.append(table_md)
                
                if text:
                    text_parts.append(text)
        
        return "\n".join(text_parts) if text_parts else None
    except Exception as e:
        print(f"pdfplumber error: {e}")
        return None

def markdown_table(table):
    """将pdfplumber表格转为Markdown"""
    if not table or not table[0]:
        return ""
    
    rows = []
    for i, row in enumerate(table):
        cells = [str(cell).strip() if cell else "" for cell in row]
        rows.append("| " + " | ".join(cells) + " |")
        if i == 0:
            rows.append("|" + "|".join(["---"] * len(cells)) + "|")
    
    return "\n".join(rows)

def read_pdf_plumber(file_path):
    """使用 pdfplumber 读取（基础版）"""
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        return "\n".join(text_parts) if text_parts else None
    except:
        return None

def read_pdf_miner(file_path):
    """使用 pdfminer.six 读取"""
    try:
        from pdfminer.high_level import extract_text
        text = extract_text(file_path)
        return text if text and len(text.strip()) > 100 else None
    except:
        return None

def read_pdf_pypdf2(file_path):
    """使用 pypdf (新版) 读取"""
    try:
        from pypdf import PdfReader
        text_parts = []
        reader = PdfReader(file_path)
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        return "\n".join(text_parts) if text_parts else None
    except:
        return None

def read_pdf(file_path):
    """读取 PDF 文档 - 多层备选，优先使用表格提取版本"""
    # 方案1: pdfplumber (带表格)
    text = read_pdf_with_tables(file_path)
    if text and len(text) > 100:
        return f"[来源: pdfplumber+tables]\n{text}"
    
    # 方案2: pdfplumber (基础)
    text = read_pdf_plumber(file_path)
    if text and len(text) > 100:
        return f"[来源: pdfplumber]\n{text}"
    
    # 方案3: pdfminer
    text = read_pdf_miner(file_path)
    if text and len(text) > 100:
        return f"[来源: pdfminer]\n{text}"
    
    # 方案4: PyPDF2
    text = read_pdf_pypdf2(file_path)
    if text and len(text) > 100:
        return f"[来源: PyPDF2]\n{text}"
    
    return "Error: 无法读取PDF文件"

def read_excel(file_path):
    """读取 Excel 文件"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        content = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            content.append(f"\n=== Sheet: {sheet_name} ===\n")
            
            # 读取表格数据，转为Markdown
            rows = list(ws.iter_rows(values_only=True))
            if rows:
                # 表头
                content.append("| " + " | ".join([str(c) if c else "" for c in rows[0]]) + " |")
                content.append("|" + "|".join(["---"] * len(rows[0])) + "|")
                # 数据行
                for row in rows[1:]:
                    if any(cell is not None for cell in row):
                        content.append("| " + " | ".join([str(cell) if cell is not None else "" for cell in row]) + " |")
        
        return "\n".join(content)
    except Exception as e:
        return f"Error reading excel: {str(e)}"

def extract_image_text_ocr(image_bytes, ext='png'):
    """使用OCR提取图片中的文字"""
    try:
        from PIL import Image
        import io
        import pytesseract
        import os
        
        # 自动配置Tesseract路径
        tess_paths = [
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
            os.path.expanduser(r'~\AppData\Local\Programs\Tesseract\tesseract.exe'),
        ]
        for path in tess_paths:
            if os.path.exists(path):
                pytesseract.pytesseract.tesseract_cmd = path
                break
        
        # 将字节转换为图片
        image = Image.open(io.BytesIO(image_bytes))
        
        # OCR识别（中文+英文）
        text = pytesseract.image_to_string(image, lang='chi_sim+eng')
        
        if text and len(text.strip()) > 5:
            return f"\n[OCR识别]: {text.strip()}\n"
        return ""
    except Exception as e:
        return f"\n[OCR失败: {str(e)[:50]}]\n"

def read_pptx(file_path, use_ocr=False):
    """
    读取 PPT 文件 - 支持文本和图片
    use_ocr: 是否使用OCR提取图片文字（需要安装pytesseract）
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        import os
        import io
        
        prs = Presentation(file_path)
        content = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = [f"\n=== Slide {i} ===\n"]
            
            # 提取文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # 提取表格文本
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        if any(row_text):
                            slide_text.append(" | ".join(row_text))
                
                # 提取图片
                if hasattr(shape, 'image'):
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        image_ext = image.ext
                        
                        img_placeholder = f"[图片{i}_{len(content)}.{image_ext}]"
                        
                        # 可选：OCR提取
                        if use_ocr:
                            ocr_text = extract_image_text_ocr(image_bytes, image_ext)
                            slide_text.append(img_placeholder + ocr_text)
                        else:
                            # 图片名称信息
                            if hasattr(image, 'name') and image.name:
                                img_placeholder += f" ({image.name})"
                            slide_text.append(img_placeholder)
                        
                    except Exception as img_err:
                        slide_text.append(f"[图片{i}_{len(content)}]")
            
            content.append("\n".join(slide_text))
        
        return "\n".join(content)
        
    except Exception as e:
        return f"Error reading pptx: {str(e)}"

def read_txt(file_path):
    """读取文本文件"""
    encodings = ['utf-8', 'gbk', 'gb2312', 'big5', 'latin1']
    for enc in encodings:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                return f.read()
        except:
            continue
    return "Error: 无法读取文本文件"

def is_heading_line(line):
    """判断是否为标题行"""
    # Markdown标题
    if line.strip().startswith('#'):
        return True
    # 常见标题模式：大写、特殊符号结尾、较短
    line = line.strip()
    if len(line) < 80 and len(line) > 2:
        # 全大写
        if line.isupper() and len(line) > 4:
            return True
        # 以特定符号结尾
        if line.endswith(':') or line.endswith('：') or line.endswith('.'):
            # 检查是否像标题（不含句子常见词）
            if not any(w in line.lower() for w in ['是', '的', '了', '在', '和', 'the', 'is', 'are', 'and']):
                return True
    return False

def semantic_chunk_text(text, chunk_size=800, overlap=100):
    """
    语义切片 - 优先按段落或标题分块
    避免切断核心知识点
    """
    if not text:
        return []
    
    chunks = []
    lines = text.split('\n')
    
    current_chunk = []
    current_size = 0
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        line_size = len(line)
        
        # 如果单行超长，强制分块
        if line_size > chunk_size:
            if current_chunk:
                chunks.append("\n".join(current_chunk))
                current_chunk = current_chunk[-3:] if len(current_chunk) > 3 else []
                current_size = sum(len(l) for l in current_chunk)
            
            # 拆分长行
            sub_chunks = []
            sub_start = 0
            while sub_start < len(line):
                sub_chunks.append(line[sub_start:sub_start + chunk_size])
                sub_start += chunk_size - overlap
            
            for sc in sub_chunks:
                chunks.append(sc)
            continue
        
        # 检查是否应该在这里分块（标题处）
        if is_heading_line(line) and current_chunk and current_size > chunk_size * 0.3:
            # 保存当前块，开始新块
            chunks.append("\n".join(current_chunk))
            current_chunk = []
            current_size = 0
        
        # 检查是否超过块大小
        if current_size + line_size > chunk_size and current_chunk:
            chunks.append("\n".join(current_chunk))
            # 保留最后几行作为overlap
            overlap_lines = current_chunk[-2:] if len(current_chunk) > 2 else current_chunk
            current_chunk = overlap_lines
            current_size = sum(len(l) for l in current_chunk)
        
        current_chunk.append(line)
        current_size += line_size + 1
    
    # 保存最后的块
    if current_chunk:
        chunks.append("\n".join(current_chunk))
    
    return chunks

def chunk_text(text, chunk_size=500, overlap=50):
    """兼容旧接口：简单分块"""
    return semantic_chunk_text(text, chunk_size, overlap)

def read_file(file_path):
    """根据文件扩展名选择读取方法"""
    if not os.path.exists(file_path):
        return f"Error: 文件不存在: {file_path}"
    
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.docx':
        return read_docx(file_path)
    elif ext == '.pdf':
        return read_pdf(file_path)
    elif ext in ['.xlsx', '.xls']:
        return read_excel(file_path)
    elif ext == '.pptx':
        return read_pptx(file_path)
    elif ext in ['.txt', '.md', '.markdown']:
        return read_txt(file_path)
    else:
        return f"Error: 不支持的文件格式: {ext}"

if __name__ == "__main__":
    if len(sys.argv) > 1:
        content = read_file(sys.argv[1])
        print(content[:1000])
    else:
        print("Usage: python read_file.py <file_path>")
